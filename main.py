import os
import tempfile
import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import contextily as ctx
import requests
from io import BytesIO
import docx
from docx.shared import Inches, Pt
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import calendar
from datetime import datetime
import math
import xyzservices.providers as xyz
from matplotlib.colors import LinearSegmentedColormap
from matplotlib.ticker import NullFormatter

# Initialize session state
if 'generated_assets' not in st.session_state:
    st.session_state.generated_assets = {}
if 'use_native_tables' not in st.session_state:
    st.session_state.use_native_tables = False

# Page configuration
st.set_page_config(layout="wide", page_title="Data Visualization and Reporting Tool")
st.title("Data Visualization and Reporting Tool")

# Helper functions for image insertion
def insert_image(doc, placeholder, asset_key):
    """Insert image into DOCX document at placeholder location"""
    if f"{{{asset_key}}}" in placeholder.text:
        placeholder.text = placeholder.text.replace(f"{{{asset_key}}}", "")
    
    buf = st.session_state.generated_assets[asset_key]
    if isinstance(buf, BytesIO):
        buf.seek(0)
    else:
        buf = BytesIO()
        st.session_state.generated_assets[asset_key].savefig(
            buf, format='png', dpi=150, bbox_inches='tight'
        )
        buf.seek(0)
    
    run = placeholder.add_run()
    run.add_picture(buf, width=Inches(6))
    buf.seek(0)

def insert_image_pptx(slide, placeholder_shape, asset_key):
    """Insert image into PPTX slide at placeholder location"""
    buf = st.session_state.generated_assets[asset_key]
    if isinstance(buf, BytesIO):
        buf.seek(0)
    else:
        buf = BytesIO()
        st.session_state.generated_assets[asset_key].savefig(
            buf, format='png', dpi=150, bbox_inches='tight'
        )
        buf.seek(0)
    
    left = placeholder_shape.left
    top = placeholder_shape.top
    width = placeholder_shape.width
    height = placeholder_shape.height
    
    slide.shapes.add_picture(
        buf, left, top, width, height
    )
    buf.seek(0)
    
    sp = placeholder_shape._element
    sp.getparent().remove(sp)

# Tab structure
tab1, tab2 = st.tabs(["Data Visualization", "Report Generation"])

def is_valid_coordinate(value):
    """Check if a value is a valid geographic coordinate"""
    try:
        num = float(value)
        return not math.isnan(num)
    except (TypeError, ValueError):
        return False

def get_climate_data(latitude, longitude):
    """Fetch daily climate data from Open-Meteo API and aggregate to monthly averages"""
    try:
        url = "https://archive-api.open-meteo.com/v1/archive"
        params = {
            "latitude": latitude,
            "longitude": longitude,
            "start_date": "2020-01-01",
            "end_date": "2023-12-31",
            "daily": ["rain_sum", "snow_depth_max"],
            "daily_snow_depth_max_unit": "cm",
            "timezone": "UTC"
        }
        response = requests.get(url, params=params)
        response.raise_for_status()
        data = response.json()
        
        if "error" in data:
            return None, None, None
        
        daily_data = data["daily"]
        df = pd.DataFrame({
            "date": pd.to_datetime(daily_data["time"]),
            "rain_sum_mm": daily_data["rain_sum"],
            "snow_depth_max_cm": daily_data["snow_depth_max"]
        })
        
        df['month'] = df['date'].dt.month
        monthly_avg = df.groupby('month').agg({
            'rain_sum_mm': 'mean',
            'snow_depth_max_cm': 'mean'
        }).reindex(range(1, 13), fill_value=0)
        
        rainfall = monthly_avg['rain_sum_mm'].tolist()
        snow_depth = monthly_avg['snow_depth_max_cm'].tolist()
        return rainfall, snow_depth, "Open-Meteo"
    
    except Exception:
        return None, None, None

def create_table_image(df):
    """Create an image of a DataFrame as a table with compact rows and no padding"""
    df = df.copy()
    for col in df.columns:
        if df[col].dtype in ['float64', 'float32']:
            df[col] = df[col].round(2).apply(lambda x: f"{x:.2f}")
    
    n_rows, n_cols = df.shape
    fig_width = max(n_cols * 0.7, 6)
    fig_height = max(n_rows * 0.18, 2)
    
    fig = plt.figure(figsize=(fig_width, fig_height), facecolor='none', dpi=300)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis('off')
    
    table = ax.table(
        cellText=df.values,
        colLabels=df.columns,
        loc='center',
        cellLoc='center',
        colLoc='center'
    )
    
    table.auto_set_font_size(False)
    table.set_fontsize(8)
    table.scale(1.0, 1.1)
    table.auto_set_column_width([i for i in range(n_cols)])
    
    for (i, j), cell in table.get_celld().items():
        if i == 0:
            cell.set_text_props(weight='bold')
        cell.set_height(0.12)
    
    buf = BytesIO()
    plt.savefig(buf, format='png', dpi=300, bbox_inches='tight', pad_inches=0, transparent=True)
    buf.seek(0)
    plt.close(fig)
    return buf

def search_docx_for_placeholder(doc, placeholder):
    """Search for placeholder in paragraphs and table cells of a DOCX document"""
    for p in doc.paragraphs:
        if placeholder in p.text:
            return p
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if placeholder in p.text:
                        return p
    return None

def search_pptx_for_placeholder(prs, placeholder):
    """Search for placeholder in all slides' shapes and table cells of a PPTX presentation"""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            return slide, shape
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if placeholder in run.text:
                                    return slide, shape
    return None, None

def add_native_table_to_docx(doc, placeholder, df):
    """Add native table to DOCX document at placeholder location"""
    parent = placeholder._element.getparent()
    table = doc.add_table(rows=df.shape[0]+1, cols=df.shape[1])
    table.style = 'Table Grid'
    table.autofit = True
    table.allow_autofit = True
    
    hdr_cells = table.rows[0].cells
    for i, col_name in enumerate(df.columns):
        hdr_cells[i].text = str(col_name)
        for paragraph in hdr_cells[i].paragraphs:
            paragraph.alignment = WD_TABLE_ALIGNMENT.CENTER
            for run in paragraph.runs:
                run.bold = True
                run.font.size = Pt(9)
    
    for i in range(df.shape[0]):
        row_cells = table.rows[i+1].cells
        for j in range(df.shape[1]):
            row_cells[j].text = str(df.iat[i, j])
            for paragraph in row_cells[j].paragraphs:
                paragraph.alignment = WD_TABLE_ALIGNMENT.CENTER
                paragraph.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(8)
    
    table_element = table._element
    parent.replace(placeholder._element, table_element)

def add_native_table_to_pptx(slide, placeholder_shape, df):
    """Add native table to PPTX slide at placeholder location"""
    left = placeholder_shape.left
    top = placeholder_shape.top
    width = placeholder_shape.width
    height = placeholder_shape.height
    
    rows, cols = df.shape[0]+1, df.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for j in range(cols):
        cell = table_shape.cell(0, j)
        cell.text = str(df.columns[j])
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(10)
    
    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            cell = table_shape.cell(i+1, j)
            cell.text = str(df.iat[i, j])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(9)

def replace_text_in_docx(doc, replacements):
    """Replace placeholders in DOCX paragraphs and table cells"""
    replaced_count = 0
    for p in doc.paragraphs:
        for key, value in replacements.items():
            if f"{{{key}}}" in p.text:
                p.text = p.text.replace(f"{{{key}}}", value)
                replaced_count += 1
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for key, value in replacements.items():
                        if f"{{{key}}}" in p.text:
                            p.text = p.text.replace(f"{{{key}}}", value)
                            replaced_count += 1
    return replaced_count

def replace_text_in_pptx(prs, replacements):
    """Replace placeholders in PPTX shapes and table cells"""
    replaced_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if f"{{{key}}}" in run.text:
                                run.text = run.text.replace(f"{{{key}}}", value)
                                replaced_count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for key, value in replacements.items():
                                    if f"{{{key}}}" in run.text:
                                        run.text = run.text.replace(f"{{{key}}}", value)
                                        replaced_count += 1
    return replaced_count

# Main tab definitions
with tab1:
    st.header("Data Visualization")
    excel_file = st.file_uploader("Upload Excel File", type=["xlsx", "xls"])
    
    nasa_api_key = st.sidebar.text_input("NASA API Key (optional, not used)", type="password")
    
    if excel_file:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
            tmp.write(excel_file.read())
            tmp_path = tmp.name

        try:
            recap_df = pd.read_excel(tmp_path, sheet_name="RECAP", header=None)
            latitude = recap_df.iloc[11, 2]
            longitude = recap_df.iloc[12, 2]
            
            if not is_valid_coordinate(latitude) or not is_valid_coordinate(longitude):
                st.error("Invalid coordinates found in RECAP sheet. Please check cells C12 (latitude) and C13 (longitude).")
                st.stop()
        except Exception as e:
            st.error(f"Error reading RECAP sheet: {str(e)}")
            st.stop()

        st.sidebar.header("Map Configuration")
        map_type = st.sidebar.selectbox(
            "Map Type",
            [
                "OpenStreetMap",
                "CartoDB Positron",
                "CartoDB Voyager",
                "Nasa Terra",
                "Nasa Topo",
                "Esri World Imagery",
                "OpenTopoMap",
                "OpenStreetMap (DE)",
                "CyclOSM",
                "OpenStreetMap (HOT)"
            ]
        )
        zoom_level = st.sidebar.slider("Zoom Level", 5, 18, 12)
        marker_color = st.sidebar.color_picker("Marker Color", "#ff0000")
        marker_size = st.sidebar.slider("Marker Size", 10, 100, 30)
        
        st.sidebar.header("Chart Styling")
        chart_style = st.sidebar.selectbox("Seaborn Style", ["whitegrid", "darkgrid", "white", "dark", "ticks"])
        
        palette_option = st.sidebar.selectbox(
            "Color Palette",
            [
                "Shades of Blue (Light to Dark)",
                "Shades of Blue (High Contrast)",
                "Seaborn: deep",
                "Seaborn: muted",
                "Seaborn: bright",
                "Seaborn: pastel",
                "Seaborn: dark",
                "Seaborn: colorblind",
                "Viridis (High Contrast)"
            ]
        )

        if palette_option == "Shades of Blue (Light to Dark)":
            colors = [
                '#e6f0ff', '#cce5ff', '#99ccff', '#66b2ff',
                '#3399ff', '#007fff', '#0055cc', '#003d99'
            ]
            sns.set_palette(sns.color_palette(colors))
        elif palette_option == "Shades of Blue (High Contrast)":
            colors = [
                '#a6cee3', '#1f78b4', '#08519c', '#08306b',
                '#6baed6', '#2171b5', '#084594', '#023858'
            ]
            sns.set_palette(sns.color_palette(colors))
        elif "Seaborn:" in palette_option:
            palette_name = palette_option.split(": ")[1].lower()
            sns.set_palette(palette_name)
        else:
            sns.set_palette("viridis")

        sns.set_style(chart_style)

        st.subheader("Location Map")
        fig, ax = plt.subplots(figsize=(10, 8))
        ax.scatter([longitude], [latitude], c=marker_color, s=marker_size, 
                  edgecolor='white', zorder=2, label="Site Location")
        
        zoom_factor = 0.02 * (18 / zoom_level)
        xlim = (longitude - zoom_factor, longitude + zoom_factor)
        ylim = (latitude - zoom_factor, latitude + zoom_factor)
        
        if all(math.isfinite(x) for x in xlim) and all(math.isfinite(y) for y in ylim):
            ax.set_xlim(xlim)
            ax.set_ylim(ylim)
        
        ax.set_aspect('equal', adjustable='datalim')
        
        try:
            providers = {
                "OpenStreetMap": xyz.OpenStreetMap.Mapnik,
                "CartoDB Positron": xyz.CartoDB.Positron,
                "CartoDB Voyager": xyz.CartoDB.Voyager,
                "Nasa Terra": xyz.Esri.NatGeoWorldMap,
                "Nasa Topo": xyz.Esri.WorldTopoMap,
                "Esri World Imagery": xyz.Esri.WorldImagery,
                "OpenTopoMap": xyz.OpenTopoMap,
                "OpenStreetMap (DE)": xyz.OpenStreetMap.DE,
                "CyclOSM": xyz.CyclOSM,
                "OpenStreetMap (HOT)": xyz.OpenStreetMap.HOT
            }
            ctx.add_basemap(ax, source=providers[map_type], crs="EPSG:4326")
        except Exception:
            pass
        
        ax.set_title("Site Location")
        ax.legend()
        st.pyplot(fig)
        st.session_state.generated_assets['chart1'] = fig
        
        rainfall, snow_depth, data_source = get_climate_data(latitude, longitude)
        source_name = data_source or "None"
        
        if rainfall is not None and snow_depth is not None:
            st.subheader(f"Weather Data ({source_name})")
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("**Monthly Average Rainfall**")
                fig, ax = plt.subplots(figsize=(8, 4))
                months = [calendar.month_abbr[i] for i in range(1, 13)]
                sns.barplot(x=months, y=rainfall, ax=ax)
                ax.set_title("Rainfall (mm/month)")
                ax.set_xlabel("Month")
                ax.set_ylabel("Rainfall (mm)")
                ax.set_ylim(0, max(rainfall)*1.2)
                st.pyplot(fig)
                st.session_state.generated_assets['chart2'] = fig
            
            with col2:
                st.markdown("**Monthly Average Snow Depth**")
                fig, ax = plt.subplots(figsize=(8, 4))
                sns.barplot(x=months, y=snow_depth, ax=ax)
                ax.set_title("Snow Depth (cm)")
                ax.set_xlabel("Month")
                ax.set_ylabel("Snow Depth (cm)")
                ax.set_ylim(0, max(snow_depth)*1.2)
                st.pyplot(fig)
                st.session_state.generated_assets['chart3'] = fig

        st.subheader("Global Horizontal Irradiation (GHI)")
        try:
            overview_df = pd.read_excel(tmp_path, sheet_name="Overview", header=3, usecols="A:G", nrows=13)
            ghi_data = overview_df.iloc[:-1].copy()
            ghi_data.iloc[:, 1:] = ghi_data.iloc[:, 1:].apply(pd.to_numeric, errors='coerce').round(2)
            months = ghi_data.iloc[:, 0].tolist()
            data_values = ghi_data.iloc[:, 1:]
            
            fig, ax = plt.subplots(figsize=(10, 5))
            for col in data_values.columns:
                sns.lineplot(x=months, y=data_values[col], label=col, linewidth=2.5, ax=ax)
            
            ax.set_title("Global Horizontal Irradiation (kWh/mÂ²/month)")
            ax.set_xlabel("Month")
            ax.set_ylabel("Irradiation")
            ax.legend(title="Data Sources")
            st.pyplot(fig)
            st.session_state.generated_assets['chart4'] = fig
        except Exception:
            pass

        st.subheader("Monthly Average Temperature")
        try:
            temp_df = pd.read_excel(tmp_path, sheet_name="Overview", header=19, usecols="A:G", nrows=13)
            temp_data = temp_df.iloc[:-1].copy()
            temp_data.iloc[:, 1:] = temp_data.iloc[:, 1:].apply(pd.to_numeric, errors='coerce').round(2)
            months = temp_data.iloc[:, 0].tolist()
            data_values = temp_data.iloc[:, 1:]
            
            fig, ax = plt.subplots(figsize=(10, 5))
            for col in data_values.columns:
                sns.lineplot(x=months, y=data_values[col], label=col, linewidth=2.5, ax=ax)
            
            ax.set_title("Temperature (Â°C)")
            ax.set_xlabel("Month")
            ax.set_ylabel("Degrees")
            ax.legend(title="Data Sources")
            st.pyplot(fig)
            st.session_state.generated_assets['chart5'] = fig
        except Exception:
            pass

        st.subheader("Probability Scenarios")
        
        try:
            prob1_df = pd.read_excel(tmp_path, sheet_name="Probability scenarios 1", 
                                    header=25, usecols="G:O", nrows=11)
            prob1_df = prob1_df.apply(lambda x: x.round(2) if x.dtype in ['float64', 'float32'] else x)
            st.markdown("**Scenario 1**")
            st.dataframe(prob1_df)
            table1_img = create_table_image(prob1_df)
            st.session_state.generated_assets['table1_img'] = table1_img
            st.session_state.generated_assets['table1_df'] = prob1_df
        except Exception:
            pass
        
        try:
            prob2_df = pd.read_excel(tmp_path, sheet_name="Probability scenarios 2", 
                                    header=2, usecols="A:L", nrows=31)
            prob2_df = prob2_df.apply(lambda x: x.round(2) if x.dtype in ['float64', 'float32'] else x)
            st.markdown("**Scenario 2**")
            st.dataframe(prob2_df)
            table2_img = create_table_image(prob2_df)
            st.session_state.generated_assets['table2_img'] = table2_img
            st.session_state.generated_assets['table2_df'] = prob2_df
        except Exception:
            pass
        
        # GAUSSIAN DISTRIBUTION CHART
        st.subheader("Gaussian Distribution")
        try:
            gaussian_df = pd.read_excel(tmp_path, sheet_name="Gaussian distribution", header=None)
            
            # Extract line data starting from row 28
            y_data = gaussian_df.iloc[27:, 3].dropna().astype(float).tolist()  # Column D (index 3)
            x_data = gaussian_df.iloc[27:, 4].dropna().astype(float).tolist()  # Column E (index 4)
            
            # Extract the three special points from rows 20,21,22
            points = []
            # Columns I, J, K correspond to indices 8,9,10
            y_columns = [8, 9, 10]  # I, J, K
            row_indices = [19, 20, 21]  # Rows 20,21,22 (0-indexed 19,20,21)
            
            for row_idx, col_idx in zip(row_indices, y_columns):
                # Get the first non-NA value in the column starting from row 28
                col_values = gaussian_df.iloc[27:, col_idx].dropna()
                if not col_values.empty:
                    y_val = col_values.iloc[0]
                    if not pd.isna(y_val):
                        # Get x-value from column E and remove decimals
                        x_val = gaussian_df.iloc[row_idx, 4]  # Column E
                        # Format as integer by converting to int and then to string
                        x_val_str = f"{int(x_val)}"
                        label = f"{gaussian_df.iloc[row_idx, 1]}, {x_val_str}"  # B column + comma + E value (integer)
                        points.append((x_val, y_val, label))
            
            # Create the plot
            fig, ax = plt.subplots(figsize=(10, 6))
            
            # Plot the main Gaussian distribution line
            sns.lineplot(x=x_data, y=y_data, ax=ax, linewidth=2.5, color='blue')
            
            # Plot the special points with annotations
            for x, y, label in points:
                ax.scatter(x, y, color='red', s=100, zorder=5)
                ax.annotate(
                    label,
                    (x, y),
                    xytext=(10, -15),
                    textcoords='offset points',
                    arrowprops=dict(arrowstyle='->', color='red', linewidth=1.5),
                    fontsize=10,
                    fontweight='bold',
                    bbox=dict(boxstyle='round,pad=0.3', fc='white', alpha=0.8, ec='red')
                )
            
            # Configure axes as requested
            ax.set_xlabel("Energy Yield [kWh/kWp/a]", fontsize=12)
            ax.set_ylabel("")  # Remove y-axis label
            ax.yaxis.set_major_formatter(NullFormatter())  # Hide y-axis values
            ax.set_title("Gaussian Distribution of Energy Yield", fontsize=14)
            ax.grid(True, linestyle='--', alpha=0.7)
            
            st.pyplot(fig)
            st.session_state.generated_assets['chart6'] = fig
            
        except Exception as e:
            st.error(f"Error processing Gaussian distribution data: {str(e)}")
        
        recap_dict = {}
        for i in range(len(recap_df)):
            key = str(recap_df.iloc[i, 0])
            value = str(recap_df.iloc[i, 2])
            if key and key.lower() != 'nan' and not key.startswith('Unnamed'):
                recap_dict[key.strip()] = value
        st.session_state.generated_assets['recap'] = recap_dict
        
        os.unlink(tmp_path)

with tab2:
    st.header("Report Generation")
    
    if not st.session_state.generated_assets:
        st.warning("Please generate assets in the Data Visualization tab first")
        st.stop()
    
    template_file = st.file_uploader("Upload Template Document", 
                                    type=["docx", "pptx"])
    output_format = st.selectbox("Select Output Format", [".docx", ".pptx"])
    
    st.session_state.use_native_tables = st.checkbox(
        "Use native tables instead of images",
        value=st.session_state.use_native_tables
    )
    
    if template_file and st.button("Generate Report"):
        with tempfile.NamedTemporaryFile(delete=False) as tmp_tpl:
            tmp_tpl.write(template_file.read())
            tpl_path = tmp_tpl.name
        
        try:
            if output_format == ".docx":
                doc = docx.Document(tpl_path)
                
                replaced_count = replace_text_in_docx(doc, st.session_state.generated_assets['recap'])
                
                charts_inserted = 0
                tables_inserted = 0
                for asset_key in ['chart1', 'chart2', 'chart3', 'chart4', 'chart5', 'chart6', 'table1_img', 'table2_img']:
                    if asset_key in st.session_state.generated_assets:
                        placeholder = search_docx_for_placeholder(doc, f"{{{asset_key}}}")
                        
                        if not placeholder:
                            st.warning(f"Placeholder {{{asset_key}}} not found in DOCX")
                            continue
                            
                        if asset_key in ['table1_img', 'table2_img'] and st.session_state.use_native_tables:
                            table_key = asset_key.replace('_img', '_df')
                            if table_key in st.session_state.generated_assets:
                                df = st.session_state.generated_assets[table_key]
                                try:
                                    add_native_table_to_docx(doc, placeholder, df)
                                    tables_inserted += 1
                                except Exception:
                                    insert_image(doc, placeholder, asset_key)
                                    tables_inserted += 1
                            else:
                                insert_image(doc, placeholder, asset_key)
                                tables_inserted += 1
                        else:
                            insert_image(doc, placeholder, asset_key)
                            if 'chart' in asset_key:
                                charts_inserted += 1
                            else:
                                tables_inserted += 1
                
                st.info(f"Generated DOCX: {replaced_count} keywords, {charts_inserted} charts, {tables_inserted} tables")
                
                output = BytesIO()
                doc.save(output)
                output.seek(0)
                st.download_button(
                    label="Download Report",
                    data=output,
                    file_name="generated_report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )
            
            else:
                prs = Presentation(tpl_path)
                
                replaced_count = replace_text_in_pptx(prs, st.session_state.generated_assets['recap'])
                
                charts_inserted = 0
                tables_inserted = 0
                for asset_key in ['chart1', 'chart2', 'chart3', 'chart4', 'chart5', 'chart6', 'table1_img', 'table2_img']:
                    if asset_key in st.session_state.generated_assets:
                        slide, placeholder_shape = search_pptx_for_placeholder(prs, f"{{{asset_key}}}")
                        
                        if not placeholder_shape:
                            st.warning(f"Placeholder {{{asset_key}}} not found in PPTX")
                            continue
                            
                        if asset_key in ['table1_img', 'table2_img'] and st.session_state.use_native_tables:
                            table_key = asset_key.replace('_img', '_df')
                            if table_key in st.session_state.generated_assets:
                                df = st.session_state.generated_assets[table_key]
                                try:
                                    add_native_table_to_pptx(slide, placeholder_shape, df)
                                    tables_inserted += 1
                                except Exception:
                                    insert_image_pptx(slide, placeholder_shape, asset_key)
                                    tables_inserted += 1
                            else:
                                insert_image_pptx(slide, placeholder_shape, asset_key)
                                tables_inserted += 1
                        else:
                            insert_image_pptx(slide, placeholder_shape, asset_key)
                            if 'chart' in asset_key:
                                charts_inserted += 1
                            else:
                                tables_inserted += 1
                
                st.info(f"Generated PPTX: {replaced_count} keywords, {charts_inserted} charts, {tables_inserted} tables")
                
                output = BytesIO()
                prs.save(output)
                output.seek(0)
                st.download_button(
                    label="Download Presentation",
                    data=output,
                    file_name="generated_report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        
        except Exception as e:
            st.error(f"Error generating report: {str(e)}")
        finally:
            os.unlink(tpl_path)
